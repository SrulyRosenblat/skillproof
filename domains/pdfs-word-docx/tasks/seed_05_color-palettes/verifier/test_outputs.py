"""Deterministic checks for /workspace/deck.pptx.

Slide 1 must be styled with Anthropic's real official brand colors (light
background, dark title text, three accent swatches in official
primary/secondary/tertiary accent order). Slide 2 must be styled with the
"Coral Energy" presentation color palette (primary background, secondary
content block, accent title text). Both palettes are made up of specific,
non-guessable hex values, so a plausible-but-wrong palette (e.g. generic navy
or orange) is deterministically caught.
"""

from pathlib import Path

import pytest
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_FILL_TYPE

DECK_PATH = Path("/workspace/deck.pptx")

BRAND_LIGHT = RGBColor(0xFA, 0xF9, 0xF5)
BRAND_DARK = RGBColor(0x14, 0x14, 0x13)
BRAND_ACCENT_PRIMARY = RGBColor(0xD9, 0x77, 0x57)
BRAND_ACCENT_SECONDARY = RGBColor(0x6A, 0x9B, 0xCC)
BRAND_ACCENT_TERTIARY = RGBColor(0x78, 0x8C, 0x5D)

CORAL_PRIMARY = RGBColor(0xF9, 0x61, 0x67)
CORAL_SECONDARY = RGBColor(0xF9, 0xE7, 0x95)
CORAL_ACCENT = RGBColor(0x2F, 0x3C, 0x7E)

SLIDE1_TITLE = "About Anthropic"
SLIDE2_TITLE = "Our Partnership Vision"


def background_rgb(slide):
    fill = slide.background.fill
    assert fill.type == MSO_FILL_TYPE.SOLID, "slide background must be a solid fill"
    return fill.fore_color.rgb


def shape_by_name(slide, name):
    for shape in slide.shapes:
        if shape.name == name:
            return shape
    raise AssertionError(f"no shape named {name!r} found on slide")


def shape_with_exact_text(slide, text):
    for shape in slide.shapes:
        if shape.has_text_frame and shape.text_frame.text == text:
            return shape
    raise AssertionError(f"no shape with exact text {text!r} found on slide")


def first_run_font_color(shape):
    for para in shape.text_frame.paragraphs:
        for run in para.runs:
            if run.text:
                return run.font.color.rgb
    raise AssertionError("shape has no non-empty run to read a font color from")


@pytest.fixture(scope="module")
def prs():
    assert DECK_PATH.exists(), f"{DECK_PATH} was not created"
    return Presentation(DECK_PATH)


def test_deck_has_exactly_two_slides(prs):
    assert len(prs.slides) == 2


def test_slide1_uses_anthropic_brand_background_and_title_color(prs):
    slide1 = prs.slides[0]
    assert background_rgb(slide1) == BRAND_LIGHT, (
        "slide 1 background must be Anthropic's official light brand color"
    )
    title_shape = shape_with_exact_text(slide1, SLIDE1_TITLE)
    assert first_run_font_color(title_shape) == BRAND_DARK, (
        "slide 1 title text must use Anthropic's official dark brand color"
    )


def test_slide1_accent_swatches_match_official_order_and_colors(prs):
    slide1 = prs.slides[0]
    one = shape_by_name(slide1, "AccentOne")
    two = shape_by_name(slide1, "AccentTwo")
    three = shape_by_name(slide1, "AccentThree")

    assert one.fill.fore_color.rgb == BRAND_ACCENT_PRIMARY, (
        "AccentOne must be Anthropic's official primary accent color"
    )
    assert two.fill.fore_color.rgb == BRAND_ACCENT_SECONDARY, (
        "AccentTwo must be Anthropic's official secondary accent color"
    )
    assert three.fill.fore_color.rgb == BRAND_ACCENT_TERTIARY, (
        "AccentThree must be Anthropic's official tertiary accent color"
    )
    assert one.left < two.left < three.left, (
        "accent swatches must be arranged left to right as AccentOne, AccentTwo, AccentThree"
    )


def test_slide2_uses_coral_energy_background_and_title_color(prs):
    slide2 = prs.slides[1]
    assert background_rgb(slide2) == CORAL_PRIMARY, (
        "slide 2 background must be the 'Coral Energy' palette's primary color"
    )
    title_shape = shape_with_exact_text(slide2, SLIDE2_TITLE)
    assert first_run_font_color(title_shape) == CORAL_ACCENT, (
        "slide 2 title text must use the 'Coral Energy' palette's accent color"
    )


def test_slide2_secondary_block_matches_coral_energy_secondary(prs):
    slide2 = prs.slides[1]
    block = shape_by_name(slide2, "SecondaryBlock")
    assert block.fill.fore_color.rgb == CORAL_SECONDARY, (
        "SecondaryBlock must be the 'Coral Energy' palette's secondary color"
    )
